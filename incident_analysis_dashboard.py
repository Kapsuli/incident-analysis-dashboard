import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import numpy as np
from datetime import datetime, timedelta
import calendar
import io
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import plotly.io as pio

# Sivun konfiguraatio
st.set_page_config(
    page_title="Hälytysten Analyysihallinta",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Suomenkieliset nimet
FINNISH_MONTHS = [
    "Tammikuu", "Helmikuu", "Maaliskuu", "Huhtikuu", "Toukokuu", "Kesäkuu",
    "Heinäkuu", "Elokuu", "Syyskuu", "Lokakuu", "Marraskuu", "Joulukuu"
]

FINNISH_WEEKDAYS = ["Ma", "Ti", "Ke", "To", "Pe", "La", "Su"]
FINNISH_WEEKDAYS_LONG = ["Maanantai", "Tiistai", "Keskiviikko", "Torstai", "Perjantai", "Lauantai", "Sunnuntai"]

def get_worker_count(hour):
    """Laske työntekijämäärä tunnin perusteella"""
    workers = 0
    
    # Yövuoro 19:15-07:15 (2 henkilöä)
    if hour >= 19 or hour < 7:
        workers += 2
    
    # Aamuvuoro 07:00-17:00 (3 henkilöä)
    if hour >= 7 and hour < 17:
        workers += 3
    
    # Iltavuorot (portaittain)
    if hour >= 9 and hour < 19: workers += 1  # 09:15-19:15
    if hour >= 10 and hour < 20: workers += 1  # 10:00-20:00
    if hour >= 11 and hour < 21: workers += 1  # 11:00-21:00
    if hour >= 13 and hour < 23: workers += 1  # 13:00-23:00
    
    return workers

def get_finnish_weekday(date_obj):
    """Palauttaa suomenkielisen viikonpäivän nimen"""
    if pd.isna(date_obj):
        return "Tuntematon"
    weekday_num = date_obj.weekday()  # 0=Maanantai, 6=Sunnuntai
    return FINNISH_WEEKDAYS_LONG[weekday_num]

def get_finnish_month_name(month_num):
    """Palauttaa suomenkielisen kuukauden nimen"""
    return FINNISH_MONTHS[month_num - 1]

def create_calendar_view(daily_stats):
    """Luo kalenterinäkymä päivittäisistä tilastoista"""
    if len(daily_stats) == 0:
        return None
    
    # Muunna päivämäärät datetime-objekteiksi
    daily_stats = daily_stats.copy()
    daily_stats['date_obj'] = pd.to_datetime(daily_stats['date'])
    
    # Määritä kuukausi ja vuosi
    first_date = daily_stats['date_obj'].min()
    last_date = daily_stats['date_obj'].max()
    
    # Jos kaikki data on samalta kuukaudelta
    if first_date.month == last_date.month and first_date.year == last_date.year:
        month = first_date.month
        year = first_date.year
    else:
        # Käytä ensimmäistä kuukautta
        month = first_date.month
        year = first_date.year
    
    # Luo kuukauden kalenteri
    cal = calendar.monthcalendar(year, month)
    
    # Luo HTML-taulukko
    month_name = get_finnish_month_name(month)
    
    # Aloita kalenteri-HTML
    calendar_html = f"""
    <div style="margin: 20px 0; font-family: Arial, sans-serif;">
        <h3 style="text-align: center; margin-bottom: 20px; color: #1f77b4; font-size: 24px;">
            📅 {month_name} {year}
        </h3>
        <div style="text-align: center; margin-bottom: 15px; font-size: 14px;">
            <span style="color: #666;">Rauhallisin: </span>
            <span style="background-color: #d4edda; padding: 4px 12px; border-radius: 5px; font-weight: bold;">
                {daily_stats.loc[daily_stats['total_incidents'].idxmin(), 'day']:.0f}. päivä ({daily_stats['total_incidents'].min():.0f} inc)
            </span>
        </div>
        <table style="width: 100%; border-collapse: collapse; margin: 0 auto; max-width: 1000px; box-shadow: 0 2px 8px rgba(0,0,0,0.1);">
            <thead>
                <tr style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white;">
    """
    
    # Viikonpäivien otsikot
    for day_name in FINNISH_WEEKDAYS:
        calendar_html += f'<th style="padding: 15px 10px; text-align: center; font-weight: bold; font-size: 16px;">{day_name}</th>'
    
    calendar_html += """
                </tr>
            </thead>
            <tbody>
    """
    
    # Kalenterin rivit
    for week in cal:
        calendar_html += '<tr>'
        for day in week:
            if day == 0:
                # Tyhjä päivä
                calendar_html += '<td style="padding: 20px; border: 1px solid #e0e0e0; background-color: #f8f9fa; height: 85px;"></td>'
            else:
                # Etsi päivän data
                day_data = daily_stats[daily_stats['date_obj'].dt.day == day]
                
                if len(day_data) > 0:
                    row = day_data.iloc[0]
                    
                    # Määritä väri tavoitteiden perusteella
                    if row['day_target_met'] and row['night_target_met']:
                        bg_color = "#d4edda"  # Vihreä - molemmat tavoitteet täytetty
                        border_color = "#28a745"
                        border_width = "3px"
                    elif row['day_target_met'] or row['night_target_met']:
                        bg_color = "#fff3cd"  # Keltainen - yksi tavoite täytetty
                        border_color = "#ffc107"
                        border_width = "2px"
                    else:
                        bg_color = "#f8d7da"  # Punainen - kumpikaan tavoite ei täytetty
                        border_color = "#dc3545"
                        border_width = "2px"
                    
                    # Määritä P: ja Y: tekstien värit tavoitteiden mukaan
                    day_text_color = "#28a745" if row['day_target_met'] else "#dc3545"  # Vihreä jos tavoite täyttyy, muuten punainen
                    night_text_color = "#28a745" if row['night_target_met'] else "#dc3545"  # Vihreä jos tavoite täyttyy, muuten punainen
                    
                    calendar_html += f"""
                    <td style="padding: 10px; border: {border_width} solid {border_color}; background-color: {bg_color}; vertical-align: top; height: 85px; position: relative; transition: all 0.3s ease;">
                        <div style="font-weight: bold; font-size: 18px; margin-bottom: 6px; color: #333;">{day}</div>
                        <div style="font-size: 11px; line-height: 1.3;">
                            <div style="color: {day_text_color}; font-weight: bold; margin-bottom: 1px;">P: {row['day_shift_avg']:.2f}</div>
                            <div style="color: {night_text_color}; font-weight: bold; margin-bottom: 1px;">Y: {row['night_shift_avg']:.2f}</div>
                            <div style="color: #666; font-size: 10px; background-color: rgba(255,255,255,0.7); padding: 1px 3px; border-radius: 3px; display: inline-block;">{row['total_incidents']:.0f} inc</div>
                        </div>
                    </td>
                    """
                else:
                    # Ei dataa tälle päivälle
                    calendar_html += f"""
                    <td style="padding: 10px; border: 1px solid #dee2e6; background-color: #ffffff; vertical-align: top; height: 85px;">
                        <div style="font-weight: bold; color: #999; font-size: 16px; margin-bottom: 4px;">{day}</div>
                        <div style="font-size: 10px; color: #ccc; font-style: italic;">Ei dataa</div>
                    </td>
                    """
        calendar_html += '</tr>'
    
    calendar_html += """
            </tbody>
        </table>
        <div style="margin-top: 20px; padding: 15px; background-color: #f8f9fa; border-radius: 8px;">
            <div style="text-align: center; font-size: 13px; color: #666; margin-bottom: 10px;">
                <strong>Selite:</strong>
            </div>
            <div style="display: flex; justify-content: center; gap: 25px; flex-wrap: wrap;">
                <span style="display: flex; align-items: center; gap: 8px;">
                    <span style="display: inline-block; width: 16px; height: 16px; background-color: #d4edda; border: 2px solid #28a745; border-radius: 3px;"></span>
                    <span style="font-size: 13px; font-weight: 500;">Molemmat tavoitteet täytetty</span>
                </span>
                <span style="display: flex; align-items: center; gap: 8px;">
                    <span style="display: inline-block; width: 16px; height: 16px; background-color: #fff3cd; border: 2px solid #ffc107; border-radius: 3px;"></span>
                    <span style="font-size: 13px; font-weight: 500;">Yksi tavoite täytetty</span>
                </span>
                <span style="display: flex; align-items: center; gap: 8px;">
                    <span style="display: inline-block; width: 16px; height: 16px; background-color: #f8d7da; border: 2px solid #dc3545; border-radius: 3px;"></span>
                    <span style="font-size: 13px; font-weight: 500;">Kumpikaan tavoite ei täytetty</span>
                </span>
                <span style="display: flex; align-items: center; gap: 8px;">
                    <span style="display: inline-block; width: 16px; height: 16px; background-color: #ffffff; border: 2px solid #dee2e6; border-radius: 3px;"></span>
                    <span style="font-size: 13px; font-weight: 500;">Ei dataa</span>
                </span>
            </div>
            <div style="text-align: center; margin-top: 12px; font-size: 12px; color: #888;">
                <strong style="color: #28a745;">P:</strong> Päivätyöntekijät (tavoite ≥5.1) | <strong style="color: #6f42c1;">Y:</strong> Yötyöntekijät (tavoite ≥4.6) | <strong>inc:</strong> Incidentit yhteensä
            </div>
        </div>
    </div>
    """
    
    return calendar_html

@st.cache_data
def process_data(df):
    """Käsittele Excel-data analyysiin"""
    try:
        # Tarkista että tarvittavat sarakkeet löytyvät
        required_columns = ['Hour', 'Incidents handled by agent']
        for col in required_columns:
            if col not in df.columns:
                st.error(f"Saraketta '{col}' ei löydy datasta. Tarkista Excel-tiedosto.")
                return None
        
        # Tee kopio datasta
        df_clean = df.copy()
        
        # Muunna Hour-sarake numeroiksi
        try:
            df_clean['Hour'] = pd.to_numeric(df_clean['Hour'], errors='coerce')
        except Exception as e:
            st.error(f"Virhe Hour-sarakkeen muunnossa: {str(e)}")
            return None
        
        # Muunna Incidents-sarake numeroiksi
        try:
            df_clean['Incidents handled by agent'] = pd.to_numeric(df_clean['Incidents handled by agent'], errors='coerce')
        except Exception as e:
            st.error(f"Virhe Incidents-sarakkeen muunnossa: {str(e)}")
            return None
        
        # Suodata vain validi data
        df_clean = df_clean[
            (df_clean['Hour'].notna()) & 
            (df_clean['Incidents handled by agent'].notna()) &
            (df_clean['Hour'] >= 0) & 
            (df_clean['Hour'] <= 23)
        ].copy()
        
        if len(df_clean) == 0:
            st.error("Ei validia dataa löydetty. Tarkista että Hour-sarake sisältää numeroita 0-23 ja Incidents-sarake sisältää numeroita.")
            return None
        
        # Lisää työntekijämäärät ja laskelmat
        df_clean['workers'] = df_clean['Hour'].apply(get_worker_count)
        df_clean['incidents_per_worker'] = df_clean['Incidents handled by agent'] / df_clean['workers']
        
        # Käsittele päivämäärät
        if 'Date' in df.columns:
            try:
                if df_clean['Date'].dtype in ['int64', 'float64']:
                    df_clean['date'] = pd.to_datetime('1900-01-01') + pd.to_timedelta(df_clean['Date'] - 2, unit='D')
                else:
                    df_clean['date'] = pd.to_datetime(df_clean['Date'], errors='coerce')
                
                if df_clean['date'].isna().all():
                    df_clean['date'] = datetime.now().date()
                    df_clean['date_str'] = df_clean['date'].astype(str)
                    df_clean['day_name'] = 'Tuntematon'
                    df_clean['day'] = 1
                else:
                    df_clean['date_str'] = df_clean['date'].dt.strftime('%Y-%m-%d')
                    df_clean['day_name'] = df_clean['date'].apply(get_finnish_weekday)
                    df_clean['day'] = df_clean['date'].dt.day
            except Exception as e:
                st.warning(f"Päivämäärien käsittely epäonnistui: {str(e)}. Käytetään oletuspäivämääriä.")
                df_clean['date'] = datetime.now().date()
                df_clean['date_str'] = df_clean['date'].astype(str)
                df_clean['day_name'] = 'Tuntematon'
                df_clean['day'] = 1
        else:
            df_clean['date'] = datetime.now().date()
            df_clean['date_str'] = df_clean['date'].astype(str)
            df_clean['day_name'] = 'Tuntematon'
            df_clean['day'] = 1
        
        return df_clean
        
    except Exception as e:
        st.error(f"Virhe datan käsittelyssä: {str(e)}")
        return None

def calculate_hourly_stats(df):
    """Laske tuntikohtaiset tilastot"""
    hourly_stats = []
    
    for hour in range(24):
        hour_data = df[df['Hour'] == hour]
        if len(hour_data) > 0:
            avg_incidents = hour_data['Incidents handled by agent'].mean()
            worker_count = get_worker_count(hour)
            avg_incidents_per_worker = avg_incidents / worker_count
            
            hourly_stats.append({
                'hour': hour,
                'hour_str': f"{hour:02d}:00",
                'avg_incidents': round(avg_incidents, 2),
                'worker_count': worker_count,
                'incidents_per_worker': round(avg_incidents_per_worker, 2),
                'days_count': len(hour_data)
            })
    
    return pd.DataFrame(hourly_stats)

def calculate_daily_stats(df):
    """Laske päivittäiset tilastot"""
    daily_stats = []
    
    for date_str in df['date_str'].unique():
        day_data = df[df['date_str'] == date_str]
        
        # Jaa päivä- ja yötyöntekijöihin
        day_shift = day_data[(day_data['Hour'] >= 7) & (day_data['Hour'] < 23)]
        night_shift = day_data[(day_data['Hour'] >= 23) | (day_data['Hour'] < 7)]
        
        day_shift_avg = day_shift['incidents_per_worker'].mean() if len(day_shift) > 0 else 0
        night_shift_avg = night_shift['incidents_per_worker'].mean() if len(night_shift) > 0 else 0
        
        daily_stats.append({
            'date': date_str,
            'day_name': day_data['day_name'].iloc[0] if len(day_data) > 0 else 'Tuntematon',
            'day': day_data['day'].iloc[0] if len(day_data) > 0 else 1,
            'total_incidents': day_data['Incidents handled by agent'].sum(),
            'day_shift_avg': round(day_shift_avg, 2),
            'night_shift_avg': round(night_shift_avg, 2),
            'day_target_met': day_shift_avg >= 5.1,
            'night_target_met': night_shift_avg >= 4.6
        })
    
    return pd.DataFrame(daily_stats)

def create_combined_chart(hourly_df):
    """Luo yhdistetty kaavio paremmilla tooltip-näkymillä"""
    fig = make_subplots(
        rows=1, cols=1,
        specs=[[{"secondary_y": True}]],
        subplot_titles=["Yhdistetty analyysi"]
    )
    
    # Pylväskaavio incidenteille - parannetulla tooltip
    fig.add_trace(
        go.Bar(
            x=hourly_df['hour_str'],
            y=hourly_df['avg_incidents'],
            name='Keskimääräiset incidentit',
            marker_color='lightblue',
            hovertemplate='<b>Kelloaika:</b> %{x}<br>' +
                         '<b>Keskimääräiset incidentit:</b> %{y:.2f}<br>' +
                         '<extra></extra>'
        ),
        secondary_y=False
    )
    
    # Viivakaavio incidenteille per työntekijä - parannetulla tooltip
    fig.add_trace(
        go.Scatter(
            x=hourly_df['hour_str'],
            y=hourly_df['incidents_per_worker'],
            name='Incidentit/työntekijä',
            line=dict(color='red', width=3),
            mode='lines+markers',
            marker=dict(size=8),
            hovertemplate='<b>Kelloaika:</b> %{x}<br>' +
                         '<b>Incidentit/työntekijä:</b> %{y:.2f}<br>' +
                         '<extra></extra>'
        ),
        secondary_y=True
    )
    
    # Viivakaavio työntekijämäärille - parannetulla tooltip
    fig.add_trace(
        go.Scatter(
            x=hourly_df['hour_str'],
            y=hourly_df['worker_count'],
            name='Työntekijämäärä',
            line=dict(color='green', width=2),
            mode='lines+markers',
            marker=dict(size=6),
            hovertemplate='<b>Kelloaika:</b> %{x}<br>' +
                         '<b>Työntekijämäärä:</b> %{y}<br>' +
                         '<extra></extra>'
        ),
        secondary_y=True
    )
    
    fig.update_xaxes(title_text="Kelloaika")
    fig.update_yaxes(title_text="Incidentit", secondary_y=False)
    fig.update_yaxes(title_text="Incidentit/työntekijä & Työntekijämäärä", secondary_y=True)
    
    fig.update_layout(
        height=500, 
        showlegend=True,
        hovermode='x unified',
        hoverlabel=dict(
            bgcolor="white",
            font_size=14,
            font_family="Arial"
        )
    )
    
    return fig

def save_plotly_as_image(fig, filename, width=1200, height=600):
    """Tallenna Plotly-kuvaaja PNG-muodossa"""
    try:
        img_bytes = pio.to_image(fig, format="png", width=width, height=height, scale=2)
        return img_bytes
    except Exception as e:
        st.error(f"Virhe kuvaajan tallentamisessa: {str(e)}")
        return None

def save_calendar_as_image(calendar_html, width=1200, height=800):
    """Tallenna kalenteri HTML:stä kuvaksi (simuloitu)"""
    # Tämä on placeholder-funktio. Oikeassa toteutuksessa käytettäisiin
    # esim. selenium + webdriver tai html2image kirjastoa
    # Nyt palautetaan None, jotta PowerPoint-toiminto toimii muuten
    return None

def create_powerpoint_presentation(selected_slides, data_dict):
    """Luo PowerPoint-esitys valituista dioista"""
    try:
        # Luo uusi esitys
        prs = Presentation()
        
        # Luo kuvat etukäteen valituille dioille
        images = {}
        
        # Tallenna yhdistetty kaavio jos tarvitaan
        if "Yhdistetty kaavio" in selected_slides:
            hourly_stats = data_dict.get('hourly_stats')
            if hourly_stats is not None and len(hourly_stats) > 0:
                try:
                    fig_combined = create_combined_chart(hourly_stats)
                    img_bytes = save_plotly_as_image(fig_combined, "combined_chart", 1400, 700)
                    if img_bytes:
                        images['combined_chart'] = img_bytes
                except Exception as e:
                    st.warning(f"Ei voitu tallentaa yhdistettyä kaaviota: {str(e)}")
        
        # Tallenna tuntikohtainen analyysi kaavio jos tarvitaan
        if "Tuntikohtainen analyysi" in selected_slides:
            hourly_stats = data_dict.get('hourly_stats')
            if hourly_stats is not None and len(hourly_stats) > 0:
                try:
                    # Luo incidentit per työntekijä kaavio
                    fig_hourly = px.line(
                        hourly_stats, 
                        x='hour_str', 
                        y='incidents_per_worker',
                        title='Incidentit per työntekijä tunnissa',
                        markers=True
                    )
                    fig_hourly.add_hline(y=5.1, line_dash="dash", line_color="red", 
                                        annotation_text="Päivätyöntekijöiden tavoite (5.1)")
                    fig_hourly.add_hline(y=4.6, line_dash="dash", line_color="blue", 
                                        annotation_text="Yötyöntekijöiden tavoite (4.6)")
                    fig_hourly.update_layout(height=500)
                    
                    img_bytes = save_plotly_as_image(fig_hourly, "hourly_analysis", 1400, 700)
                    if img_bytes:
                        images['hourly_analysis'] = img_bytes
                except Exception as e:
                    st.warning(f"Ei voitu tallentaa tuntikohtaista analyysiä: {str(e)}")
        
        # Tallenna kuukausinäkymä kaavio jos tarvitaan
        if "Kuukausinäkymä" in selected_slides:
            daily_stats = data_dict.get('daily_stats')
            if daily_stats is not None and len(daily_stats) > 0:
                try:
                    # Luo päivittäinen kehitys kaavio
                    fig_daily = px.line(
                        daily_stats, 
                        x='date', 
                        y=['day_shift_avg', 'night_shift_avg'],
                        title='Päivittäinen kehitys - Tuottavuustavoitteet',
                        labels={
                            'value': 'Inc/työnt./h', 
                            'variable': 'Vuoro',
                            'date': 'Päivämäärä'
                        }
                    )
                    
                    # Muuta legendan nimet suomeksi
                    fig_daily.for_each_trace(
                        lambda t: t.update(
                            name='Päivätyöntekijät' if 'day_shift_avg' in t.name else 'Yötyöntekijät'
                        )
                    )
                    
                    fig_daily.add_hline(y=5.1, line_dash="dash", line_color="red", 
                                      annotation_text="Päivätyöntekijöiden tavoite (5.1)")
                    fig_daily.add_hline(y=4.6, line_dash="dash", line_color="blue", 
                                      annotation_text="Yötyöntekijöiden tavoite (4.6)")
                    fig_daily.update_layout(height=500)
                    
                    img_bytes = save_plotly_as_image(fig_daily, "monthly_view", 1400, 700)
                    if img_bytes:
                        images['monthly_view'] = img_bytes
                except Exception as e:
                    st.warning(f"Ei voitu tallentaa kuukausinäkymää: {str(e)}")
        
        # Luo diat käyttäen oikeita parametreja
        for slide_type in selected_slides:
            try:
                if slide_type == "Yhteenveto":
                    create_summary_slide(prs, data_dict)
                elif slide_type == "Tuottavuustavoitteet":
                    create_targets_slide(prs, data_dict)
                elif slide_type == "Tuntikohtainen analyysi":
                    create_hourly_analysis_slide(prs, data_dict, images.get('hourly_analysis'))
                elif slide_type == "Kuukausinäkymä":
                    create_monthly_view_slide(prs, data_dict, images.get('monthly_view'))
                elif slide_type == "Suositukset":
                    create_recommendations_slide(prs, data_dict)
                elif slide_type == "Yhdistetty kaavio":
                    create_combined_chart_slide(prs, data_dict, images.get('combined_chart'))
            except TypeError as te:
                if "positional arguments" in str(te):
                    st.warning(f"⚠️ Ohitetaan dia '{slide_type}' - parametrivirhe: {str(te)}")
                    # Yritä ilman kuvaparametria
                    try:
                        if slide_type == "Yhdistetty kaavio":
                            create_combined_chart_slide_fallback(prs, data_dict)
                        elif slide_type == "Tuntikohtainen analyysi":
                            create_hourly_analysis_slide_text_only(prs, data_dict)
                        elif slide_type == "Kuukausinäkymä":
                            create_monthly_view_slide_text_only(prs, data_dict)
                    except:
                        st.error(f"❌ Dia '{slide_type}' ohitettiin kokonaan")
                else:
                    raise te
        
        return prs
    
    except Exception as e:
        st.error(f"Virhe PowerPoint-esityksen luonnissa: {str(e)}")
        return None

def create_summary_slide(prs, data_dict):
    """Luo yhteenveto-dia"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Otsikko
    title = slide.shapes.title
    title.text = "📊 Hälytysten Analyysi - Yhteenveto"
    
    # Sisältö
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Lisää yhteenveto-tiedot
    processed_df = data_dict.get('processed_df')
    daily_stats = data_dict.get('daily_stats')
    day_avg = data_dict.get('day_avg', 0)
    night_avg = data_dict.get('night_avg', 0)
    
    if processed_df is not None and len(processed_df) > 0:
        total_incidents = processed_df['Incidents handled by agent'].sum()
        analysis_period = f"{processed_df['date_str'].min()} - {processed_df['date_str'].max()}"
        total_days = len(daily_stats) if daily_stats is not None else 1
        
        p = tf.paragraphs[0]
        p.text = f"Analysoitu ajanjakso: {analysis_period}"
        p.font.size = Pt(16)
        p.font.bold = True
        
        # Lisää tilastoja
        stats_text = f"""
• Yhteensä incidenttejä: {total_incidents:.0f}
• Analysoitu päiviä: {total_days}
• Keskimäärin incidenttejä/päivä: {total_incidents/total_days:.1f}

Tuottavuustavoitteiden tulokset:
• Päivätyöntekijät (07-23): {day_avg:.2f} inc/työnt./h
  - Tavoite: ≥5.1 inc/työnt./h
  - Tulos: {'✅ SAAVUTETTU' if day_avg >= 5.1 else '❌ EI SAAVUTETTU'}
  
• Yötyöntekijät (23-07): {night_avg:.2f} inc/työnt./h
  - Tavoite: ≥4.6 inc/työnt./h
  - Tulos: {'✅ SAAVUTETTU' if night_avg >= 4.6 else '❌ EI SAAVUTETTU'}
"""
        
        for line in stats_text.strip().split('\n'):
            if line.strip():
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                if '✅' in line or '❌' in line:
                    p.font.bold = True

def create_targets_slide(prs, data_dict):
    """Luo tuottavuustavoitteet-dia"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🎯 Tuottavuustavoitteiden Analyysi"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    day_avg = data_dict.get('day_avg', 0)
    night_avg = data_dict.get('night_avg', 0)
    
    # Päivätyöntekijät
    p = tf.paragraphs[0]
    p.text = "🌅 PÄIVÄTYÖNTEKIJÄT (07:00-23:00)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 0) if day_avg >= 5.1 else RGBColor(200, 0, 0)
    
    p = tf.add_paragraph()
    p.text = f"Keskiarvo: {day_avg:.2f} inc/työnt./h"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = f"Tavoite: ≥5.1 inc/työnt./h"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = f"Ero tavoitteeseen: {day_avg - 5.1:+.2f}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 100, 0) if day_avg >= 5.1 else RGBColor(200, 0, 0)
    
    # Tyhjä rivi
    tf.add_paragraph()
    
    # Yötyöntekijät
    p = tf.add_paragraph()
    p.text = "🌙 YÖTYÖNTEKIJÄT (23:00-07:00)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 0) if night_avg >= 4.6 else RGBColor(200, 0, 0)
    
    p = tf.add_paragraph()
    p.text = f"Keskiarvo: {night_avg:.2f} inc/työnt./h"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = f"Tavoite: ≥4.6 inc/työnt./h"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = f"Ero tavoitteeseen: {night_avg - 4.6:+.2f}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 100, 0) if night_avg >= 4.6 else RGBColor(200, 0, 0)

def create_hourly_analysis_slide(prs, data_dict):
    """Luo tuntikohtainen analyysi -dia"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📈 Tuntikohtainen Analyysi"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    hourly_stats = data_dict.get('hourly_stats')
    
    if hourly_stats is not None and len(hourly_stats) > 0:
        # Etsi huipputunnit
        max_incidents_hour = hourly_stats.loc[hourly_stats['avg_incidents'].idxmax()]
        min_incidents_hour = hourly_stats.loc[hourly_stats['avg_incidents'].idxmin()]
        max_efficiency_hour = hourly_stats.loc[hourly_stats['incidents_per_worker'].idxmax()]
        min_efficiency_hour = hourly_stats.loc[hourly_stats['incidents_per_worker'].idxmin()]
        
        p = tf.paragraphs[0]
        p.text = "KESKEISET HAVAINNOT:"
        p.font.size = Pt(16)
        p.font.bold = True
        
        findings = [
            f"• Kiireisin tunti: {max_incidents_hour['hour_str']} ({max_incidents_hour['avg_incidents']:.1f} inc)",
            f"• Rauhallisim tunti: {min_incidents_hour['hour_str']} ({min_incidents_hour['avg_incidents']:.1f} inc)",
            f"• Tehokkain tunti: {max_efficiency_hour['hour_str']} ({max_efficiency_hour['incidents_per_worker']:.2f} inc/työnt.)",
            f"• Vähiten tehokas: {min_efficiency_hour['hour_str']} ({min_efficiency_hour['incidents_per_worker']:.2f} inc/työnt.)",
            "",
            "VUOROJÄRJESTELY:",
            "• Yövuoro (19:15-07:15): 2 henkilöä",
            "• Aamuvuoro (07:00-17:00): 3 henkilöä", 
            "• Iltavuorot (portaittain): 1-4 henkilöä lisää"
        ]
        
        for finding in findings:
            if finding.strip():
                p = tf.add_paragraph()
                p.text = finding
                p.font.size = Pt(14)
                if finding.startswith("VUOROJÄRJESTELY"):
                    p.font.bold = True
                    p.font.size = Pt(15)

def create_monthly_view_slide(prs, data_dict):
    """Luo kuukausinäkymä-dia"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📅 Kuukausinäkymä"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    daily_stats = data_dict.get('daily_stats')
    
    if daily_stats is not None and len(daily_stats) > 0:
        # Kuukauden tilastot
        day_target_met = len(daily_stats[daily_stats['day_target_met']])
        night_target_met = len(daily_stats[daily_stats['night_target_met']])
        total_days = len(daily_stats)
        
        max_day = daily_stats.loc[daily_stats['total_incidents'].idxmax()]
        min_day = daily_stats.loc[daily_stats['total_incidents'].idxmin()]
        
        p = tf.paragraphs[0]
        p.text = "KUUKAUDEN YHTEENVETO:"
        p.font.size = Pt(16)
        p.font.bold = True
        
        stats = [
            f"• Analysoitu päiviä: {total_days}",
            f"• Päivätyöntekijöiden tavoite täytetty: {day_target_met}/{total_days} päivää ({day_target_met/total_days*100:.1f}%)",
            f"• Yötyöntekijöiden tavoite täytetty: {night_target_met}/{total_days} päivää ({night_target_met/total_days*100:.1f}%)",
            "",
            f"• Kiireisin päivä: {max_day['day']:.0f}. ({max_day['day_name']}) - {max_day['total_incidents']:.0f} inc",
            f"• Rauhallisin päivä: {min_day['day']:.0f}. ({min_day['day_name']}) - {min_day['total_incidents']:.0f} inc",
            "",
            "PÄIVITTÄINEN KEHITYS:",
            f"• Päivätyöntekijöiden keskiarvo: {daily_stats['day_shift_avg'].mean():.2f}",
            f"• Yötyöntekijöiden keskiarvo: {daily_stats['night_shift_avg'].mean():.2f}",
            f"• Kokonaisincidenttien keskiarvo: {daily_stats['total_incidents'].mean():.1f}/päivä"
        ]
        
        for stat in stats:
            if stat.strip():
                p = tf.add_paragraph()
                p.text = stat
                p.font.size = Pt(14)
                if stat.startswith("PÄIVITTÄINEN KEHITYS"):
                    p.font.bold = True
                    p.font.size = Pt(15)

def create_recommendations_slide(prs, data_dict):
    """Luo suositukset-dia"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💡 Optimointisuositukset"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    hourly_stats = data_dict.get('hourly_stats')
    day_avg = data_dict.get('day_avg', 0)
    night_avg = data_dict.get('night_avg', 0)
    
    p = tf.paragraphs[0]
    p.text = "TOIMENPIDESUOSITUKSET:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    if hourly_stats is not None and len(hourly_stats) > 0:
        # Analysoi ongelmatunnit
        day_problems = hourly_stats[
            (hourly_stats['hour'] >= 7) & 
            (hourly_stats['hour'] < 23) & 
            (hourly_stats['incidents_per_worker'] < 5.1)
        ]
        
        night_problems = hourly_stats[
            ((hourly_stats['hour'] >= 23) | (hourly_stats['hour'] < 7)) & 
            (hourly_stats['incidents_per_worker'] < 4.6)
        ]
        
        recommendations = []
        
        # Kokonaisarvio
        if day_avg >= 5.1 and night_avg >= 4.6:
            recommendations.append("✅ MOLEMMAT TAVOITTEET SAAVUTETTU")
            recommendations.append("• Jatka nykyisellä strategialla")
            recommendations.append("• Seuraa trendin kehitystä")
        else:
            if day_avg < 5.1:
                recommendations.append("❌ PÄIVÄTYÖNTEKIJÄT - TOIMENPITEET TARVITAAN")
                if len(day_problems) > 0:
                    problem_hours = ", ".join([f"{row['hour_str']}" for _, row in day_problems.iterrows()])
                    recommendations.append(f"• Ongelmatunnit: {problem_hours}")
                recommendations.append("• Harkitse henkilöstön lisäämistä ongelmallisina aikoina")
                recommendations.append("• Analysoi työkuormituksen jakoa")
            
            if night_avg < 4.6:
                recommendations.append("❌ YÖTYÖNTEKIJÄT - TOIMENPITEET TARVITAAN")
                if len(night_problems) > 0:
                    problem_hours = ", ".join([f"{row['hour_str']}" for _, row in night_problems.iterrows()])
                    recommendations.append(f"• Ongelmatunnit: {problem_hours}")
                recommendations.append("• Harkitse henkilöstön lisäämistä yövuoroon")
                recommendations.append("• Tarkista yövuoron prosessit")
        
        recommendations.extend([
            "",
            "JATKUVAT PARANNUSTOIMET:",
            "• Seuraa tuntikohtaisia trendejä viikoittain",
            "• Analysoi incidenttityyppien jakautumista",
            "• Optimoi vuorosuunnittelua datan perusteella",
            "• Koulutettu henkilöstöä tehokkaampiin työmenetelmiin"
        ])
        
        for rec in recommendations:
            if rec.strip():
                p = tf.add_paragraph()
                p.text = rec
                p.font.size = Pt(14)
                if rec.startswith("✅") or rec.startswith("❌"):
                    p.font.bold = True
                    p.font.size = Pt(15)
                elif rec.startswith("JATKUVAT PARANNUSTOIMET"):
                    p.font.bold = True
                    p.font.size = Pt(15)

def create_combined_chart_slide(prs, data_dict):
    """Luo yhdistetty kaavio -dia"""
    slide_layout = prs.slide_layouts[5]  # Blank layout for chart
    slide = prs.slides.add_slide(slide_layout)
    
    # Lisää otsikko
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "📊 Yhdistetty Analyysi - Tuntikohtainen Kuormitus"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Lisää selitysteksti
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1.5))
    info_frame = info_box.text_frame
    
    hourly_stats = data_dict.get('hourly_stats')
    if hourly_stats is not None and len(hourly_stats) > 0:
        max_incidents = hourly_stats['avg_incidents'].max()
        max_hour = hourly_stats.loc[hourly_stats['avg_incidents'].idxmax(), 'hour_str']
        max_efficiency = hourly_stats['incidents_per_worker'].max()
        max_eff_hour = hourly_stats.loc[hourly_stats['incidents_per_worker'].idxmax(), 'hour_str']
        
        info_text = f"""Kaavio näyttää:
• Sininen pylväs: Keskimääräiset incidentit tunnissa (maksimi {max_incidents:.1f} klo {max_hour})
• Punainen viiva: Incidentit per työntekijä (maksimi {max_efficiency:.2f} klo {max_eff_hour})
• Vihreä viiva: Työntekijämäärä vuoroittain (2-7 henkilöä)

Tavoitteet: Päivätyöntekijät ≥5.1, Yötyöntekijät ≥4.6 inc/työnt./h"""
        
        info_para = info_frame.paragraphs[0]
        info_para.text = info_text
        info_para.font.size = Pt(14)

def create_combined_chart_slide_fallback(prs, data_dict):
    """Luo yhdistetty kaavio -dia ilman kuvaa (fallback)"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📊 Yhdistetty Analyysi - Tuntikohtainen Kuormitus"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    hourly_stats = data_dict.get('hourly_stats')
    if hourly_stats is not None and len(hourly_stats) > 0:
        max_incidents = hourly_stats['avg_incidents'].max()
        max_hour = hourly_stats.loc[hourly_stats['avg_incidents'].idxmax(), 'hour_str']
        max_efficiency = hourly_stats['incidents_per_worker'].max()
        max_eff_hour = hourly_stats.loc[hourly_stats['incidents_per_worker'].idxmax(), 'hour_str']
        
        p = tf.paragraphs[0]
        p.text = "YHDISTETTY ANALYYSI:"
        p.font.size = Pt(16)
        p.font.bold = True
        
        analysis_points = [
            f"• Kiireisin tunti: {max_hour} ({max_incidents:.1f} incidenttiä)",
            f"• Tehokkain tunti: {max_eff_hour} ({max_efficiency:.2f} inc/työnt.)",
            f"• Työntekijämäärä vaihtelee 2-7 henkilön välillä vuorojen mukaan",
            "",
            "TAVOITTEET:",
            "• Päivätyöntekijät: ≥5.1 inc/työnt./h",
            "• Yötyöntekijät: ≥4.6 inc/työnt./h",
            "",
            "HUOMIO: Kaavio ei saatavilla - tarvittaisiin kaleido-kirjasto kuvien tallennukseen"
        ]
        
        for point in analysis_points:
            if point.strip():
                p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(14)
                if point.startswith("TAVOITTEET") or point.startswith("YHDISTETTY ANALYYSI"):
                    p.font.bold = True
                    p.font.size = Pt(15)

def download_powerpoint(prs, filename="incident_analysis.pptx"):
    """Luo latauslinkki PowerPoint-tiedostolle"""
    try:
        # Tallenna muistiin
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        # Luo latauslinkki
        b64 = base64.b64encode(output.read()).decode()
        href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{filename}">💾 Lataa PowerPoint-esitys</a>'
        return href
    except Exception as e:
        st.error(f"Virhe PowerPoint-tiedoston luonnissa: {str(e)}")
        return None

def main():
    # Otsikko
    st.title("📊 Hälytysten Analyysihallinta")
    st.markdown("**Lataa Excel-tiedosto ja saa automaattinen analyysi hälytysten määrästä suhteessa työntekijöihin**")
    
    # PowerPoint-asetukset pääsivulla
    st.markdown("---")
    st.subheader("📑 PowerPoint-esityksen asetukset")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.markdown("**Valitse diat joita haluat sisällyttää esitykseen:**")
        
        slide_options = [
            "Yhteenveto",
            "Tuottavuustavoitteet", 
            "Yhdistetty kaavio",
            "Tuntikohtainen analyysi",
            "Kuukausinäkymä",
            "Suositukset"
        ]
        
        selected_slides = st.multiselect(
            "Valitse diat:",
            slide_options,
            default=["Yhteenveto", "Tuottavuustavoitteet", "Suositukset"],
            help="Valitse ne diat jotka haluat sisällyttää PowerPoint-esitykseen",
            key="slide_selector"
        )
        
        if selected_slides:
            st.success(f"✅ Valittu {len(selected_slides)} diaa: {', '.join(selected_slides)}")
        else:
            st.warning("⚠️ Valitse vähintään yksi dia luodaksesi PowerPoint-esityksen")
    
    with col2:
        st.markdown("**Diojen kuvaukset:**")
        descriptions = {
            "Yhteenveto": "📊 Pääkohdat ja kokonaistilanne",
            "Tuottavuustavoitteet": "🎯 Tavoitteiden täyttyminen",
            "Yhdistetty kaavio": "📈 Visuaalinen analyysi",
            "Tuntikohtainen analyysi": "🕐 Tuntitason tarkastelu",
            "Kuukausinäkymä": "📅 Päivittäiset trendit",
            "Suositukset": "💡 Toimenpide-ehdotukset"
        }
        
        for slide in selected_slides:
            if slide in descriptions:
                st.write(f"• {descriptions[slide]}")
    
    # Sivupalkki
    with st.sidebar:
        st.header("⚙️ Asetukset")
        
        # Tiedoston lataus
        uploaded_file = st.file_uploader(
            "Lataa Excel-tiedosto",
            type=['xlsx', 'xls'],
            help="Tiedoston tulee sisältää sarakkeet: 'Hour', 'Incidents handled by agent', ja mahdollisesti 'Date'"
        )
        
        st.markdown("---")
        
        # Vuorojen selitys
        st.subheader("🕐 Vuorojärjestely")
        st.markdown("""
        **Yövuoro** (19:15-07:15): 2 henkilöä  
        **Aamuvuoro** (07:00-17:00): 3 henkilöä  
        **Iltavuorot** (portaittain): 1-4 henkilöä
        - 09:15-19:15: +1 henkilö
        - 10:00-20:00: +1 henkilö  
        - 11:00-21:00: +1 henkilö
        - 13:00-23:00: +1 henkilö
        """)
        
        st.markdown("---")
        
        # Tuottavuustavoitteet
        st.subheader("🎯 Tuottavuustavoitteet")
        st.markdown("""
        **Päivätyöntekijät** (07-23): ≥5.1 inc/työnt./h  
        **Yötyöntekijät** (23-07): ≥4.6 inc/työnt./h
        """)
    
    # Pääsisältö
    if uploaded_file is not None:
        try:
            # Lue Excel-tiedosto
            df = pd.read_excel(uploaded_file)
            st.success(f"✅ Tiedosto ladattu! Löydettiin {len(df)} riviä dataa.")
            
            # Näytä datan otsikko
            with st.expander("📋 Näytä raakadata (ensimmäiset 10 riviä)"):
                st.dataframe(df.head(10))
                
                # Näytä sarakkeiden tietotyypit
                st.subheader("Sarakkeiden tietotyypit:")
                for col in df.columns:
                    st.write(f"- **{col}**: {df[col].dtype}")
            
            # Käsittele data
            processed_df = process_data(df)
            
            if processed_df is not None:
                st.success(f"✅ Data käsitelty onnistuneesti! {len(processed_df)} validia riviä.")
                
                # Laske tilastot
                hourly_stats = calculate_hourly_stats(processed_df)
                daily_stats = calculate_daily_stats(processed_df)
                
                # Tuottavuustavoitteiden analyysi
                day_shift_data = processed_df[(processed_df['Hour'] >= 7) & (processed_df['Hour'] < 23)]
                night_shift_data = processed_df[(processed_df['Hour'] >= 23) | (processed_df['Hour'] < 7)]
                
                day_avg = day_shift_data['incidents_per_worker'].mean() if len(day_shift_data) > 0 else 0
                night_avg = night_shift_data['incidents_per_worker'].mean() if len(night_shift_data) > 0 else 0
                
                # Luo data_dict PowerPoint-generointia varten
                data_dict = {
                    'processed_df': processed_df,
                    'hourly_stats': hourly_stats,
                    'daily_stats': daily_stats,
                    'day_avg': day_avg,
                    'night_avg': night_avg
                }
                
                # PowerPoint-latausmahdollisuus
                st.markdown("---")
                st.subheader("📑 Luo PowerPoint-esitys")
                
                if selected_slides and len(selected_slides) > 0:
                    col1, col2 = st.columns([2, 1])
                    
                    with col1:
                        st.write(f"🎯 **Valitut diat ({len(selected_slides)}):** {', '.join(selected_slides)}")
                        st.write("📄 Esitys sisältää analyysitulokset ja visualisoinnit valituista osioista")
                    
                    with col2:
                        if st.button("🔄 Luo PowerPoint-esitys", type="primary", use_container_width=True):
                            with st.spinner("🔄 Luodaan PowerPoint-esitystä..."):
                                try:
                                    ppt = create_powerpoint_presentation(selected_slides, data_dict)
                                    if ppt:
                                        download_link = download_powerpoint(ppt, f"incident_analysis_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx")
                                        if download_link:
                                            st.success("✅ PowerPoint-esitys luotu onnistuneesti!")
                                            st.markdown("---")
                                            st.markdown("### 💾 Lataa esitys:")
                                            st.markdown(download_link, unsafe_allow_html=True)
                                            st.info("💡 Klikkaa linkkiä ladataksesi PowerPoint-tiedoston. Esitys sisältää kaaviot kuvina valituista osioista.")
                                        else:
                                            st.error("❌ PowerPoint-latauslinkin luonti epäonnistui")
                                    else:
                                        st.error("❌ PowerPoint-esityksen luonti epäonnistui")
                                except ImportError:
                                    st.error("❌ PowerPoint-ominaisuus vaatii python-pptx kirjaston")
                                    st.info("📦 Lisää requirements.txt tiedostoon: `python-pptx`")
                                    if st.checkbox("🔧 Näytä Streamlit Cloud ohjeet"):
                                        st.markdown("""
                                        **Streamlit Cloud:ssa:**
                                        1. Lisää projektisi juureen `requirements.txt` tiedosto
                                        2. Sisällytä seuraavat rivit:
                                        ```
                                        streamlit
                                        pandas
                                        plotly
                                        python-pptx
                                        kaleido
                                        openpyxl
                                        ```
                                        3. Commitoi muutokset GitHubiin
                                        4. Streamlit Cloud asentaa kirjastot automaattisesti
                                        """)
                                except Exception as e:
                                    st.error(f"❌ Virhe PowerPoint-esityksen luonnissa: {str(e)}")
                                    
                                    # Streamlit Cloud -spesifiset virheilmoitukset
                                    if "streamlit" in str(e).lower() or "cloud" in str(e).lower():
                                        st.info("☁️ Streamlit Cloud -ympäristössä havaittu ongelma")
                                        st.info("💡 PowerPoint toimii ilman kuvia - yritä uudelleen")
                                    elif "kaleido" in str(e).lower():
                                        st.info("📦 Kuvien tallennus ei onnistu, mutta PowerPoint luodaan ilman kuvia")
                                        st.info("☁️ Streamlit Cloud:ssa kaleido ei aina toimi - tämä on normaalia")
                                    elif "takes 2 positional arguments" in str(e):
                                        st.info("🔧 Funktioparametrien virhe - yritä valita vähemmän diat")
                                    else:
                                        st.info("💡 Varmista että requirements.txt sisältää: python-pptx, plotly, pandas")
                                        
                                        if st.checkbox("🐛 Näytä debug-tiedot"):
                                            st.code(f"Virhe: {str(e)}")
                                            st.code(f"Virhetyyppi: {type(e).__name__}")
                                    
                                    # Tarjoa vaihtoehtoinen ratkaisu
                                    st.markdown("---")
                                    st.info("🔄 **Vaihtoehtoinen ratkaisu:** Kokeile luoda PowerPoint vähemmällä dialla tai pelkällä tekstisisällöllä")
                else:
                    st.warning("⚠️ Valitse vähintään yksi dia yllä olevasta listasta luodaksesi PowerPoint-esityksen")
                    st.info("👆 Voit valita diat sivun yläosasta PowerPoint-asetuksista")
                
                # Tulosten näyttäminen
                st.header("🎯 Tuottavuustavoitteiden tulokset")
                
                col1, col2 = st.columns(2)
                
                with col1:
                    day_status = "✅ SAAVUTETTU" if day_avg >= 5.1 else "❌ EI SAAVUTETTU"
                    day_color = "green" if day_avg >= 5.1 else "red"
                    st.markdown(f"""
                    <div style="padding: 20px; border: 2px solid {day_color}; border-radius: 10px; background-color: {'lightgreen' if day_avg >= 5.1 else 'lightcoral'};">
                        <h3>🌅 Päivätyöntekijät (07-23)</h3>
                        <p><strong>Keskiarvo:</strong> {day_avg:.2f} inc/työnt./h</p>
                        <p><strong>Tavoite:</strong> ≥5.1 inc/työnt./h</p>
                        <p><strong>Tulos:</strong> {day_status}</p>
                        <p><strong>Ero:</strong> {day_avg - 5.1:+.2f}</p>
                    </div>
                    """, unsafe_allow_html=True)
                
                with col2:
                    night_status = "✅ SAAVUTETTU" if night_avg >= 4.6 else "❌ EI SAAVUTETTU"
                    night_color = "green" if night_avg >= 4.6 else "red"
                    st.markdown(f"""
                    <div style="padding: 20px; border: 2px solid {night_color}; border-radius: 10px; background-color: {'lightgreen' if night_avg >= 4.6 else 'lightcoral'};">
                        <h3>🌙 Yötyöntekijät (23-07)</h3>
                        <p><strong>Keskiarvo:</strong> {night_avg:.2f} inc/työnt./h</p>
                        <p><strong>Tavoite:</strong> ≥4.6 inc/työnt./h</p>
                        <p><strong>Tulos:</strong> {night_status}</p>
                        <p><strong>Ero:</strong> {night_avg - 4.6:+.2f}</p>
                    </div>
                    """, unsafe_allow_html=True)
                
                # Välilehdet eri näkymille
                tab1, tab2, tab3, tab4, tab5 = st.tabs([
                    "📊 Yhdistetty näkymä", 
                    "📈 Tuntikohtainen analyysi", 
                    "📅 Kuukausinäkymä",
                    "📋 Yksityiskohtaiset tilastot",
                    "💡 Suositukset"
                ])
                
                with tab1:
                    st.subheader("Yhdistetty analyysi")
                    if len(hourly_stats) > 0:
                        try:
                            fig_combined = create_combined_chart(hourly_stats)
                            st.plotly_chart(fig_combined, use_container_width=True)
                        except Exception as e:
                            st.error(f"Virhe kaavion luonnissa: {str(e)}")
                            st.info("Näytetään data taulukkona:")
                            st.dataframe(hourly_stats)
                    else:
                        st.warning("Ei dataa kaavion piirtämiseen.")
                
                with tab2:
                    st.subheader("Tuntikohtainen analyysi")
                    
                    if len(hourly_stats) > 0:
                        # Valitse näkymä
                        chart_type = st.selectbox(
                            "Valitse näkymä:",
                            ["Incidentit/työntekijä", "Kokonaisincidentit", "Työntekijämäärät"]
                        )
                        
                        try:
                            if chart_type == "Incidentit/työntekijä":
                                fig = px.line(
                                    hourly_stats, 
                                    x='hour_str', 
                                    y='incidents_per_worker',
                                    title='Incidentit per työntekijä tunnissa',
                                    markers=True,
                                    hover_data={
                                        'hour_str': False,
                                        'incidents_per_worker': ':.2f',
                                        'worker_count': True,
                                        'avg_incidents': ':.2f'
                                    }
                                )
                                fig.update_traces(
                                    hovertemplate='<b>Kelloaika:</b> %{x}<br>' +
                                                 '<b>Incidentit/työntekijä:</b> %{y:.2f}<br>' +
                                                 '<b>Työntekijämäärä:</b> %{customdata[0]}<br>' +
                                                 '<b>Keskimääräiset incidentit:</b> %{customdata[1]:.2f}<br>' +
                                                 '<extra></extra>',
                                    customdata=hourly_stats[['worker_count', 'avg_incidents']].values
                                )
                                fig.add_hline(y=5.1, line_dash="dash", line_color="red", 
                                             annotation_text="Päivätyöntekijöiden tavoite (5.1)")
                                fig.add_hline(y=4.6, line_dash="dash", line_color="blue", 
                                             annotation_text="Yötyöntekijöiden tavoite (4.6)")
                            
                            elif chart_type == "Kokonaisincidentit":
                                fig = px.bar(
                                    hourly_stats, 
                                    x='hour_str', 
                                    y='avg_incidents',
                                    title='Keskimääräiset incidentit tunneittain',
                                    hover_data={
                                        'hour_str': False,
                                        'avg_incidents': ':.2f',
                                        'worker_count': True,
                                        'incidents_per_worker': ':.2f'
                                    }
                                )
                                fig.update_traces(
                                    hovertemplate='<b>Kelloaika:</b> %{x}<br>' +
                                                 '<b>Keskimääräiset incidentit:</b> %{y:.2f}<br>' +
                                                 '<b>Työntekijämäärä:</b> %{customdata[0]}<br>' +
                                                 '<b>Incidentit/työntekijä:</b> %{customdata[1]:.2f}<br>' +
                                                 '<extra></extra>',
                                    customdata=hourly_stats[['worker_count', 'incidents_per_worker']].values
                                )
                            
                            else:  # Työntekijämäärät
                                fig = px.bar(
                                    hourly_stats, 
                                    x='hour_str', 
                                    y='worker_count',
                                    title='Työntekijämäärät tunneittain',
                                    hover_data={
                                        'hour_str': False,
                                        'worker_count': True,
                                        'avg_incidents': ':.2f',
                                        'incidents_per_worker': ':.2f'
                                    }
                                )
                                fig.update_traces(
                                    hovertemplate='<b>Kelloaika:</b> %{x}<br>' +
                                                 '<b>Työntekijämäärä:</b> %{y}<br>' +
                                                 '<b>Keskimääräiset incidentit:</b> %{customdata[0]:.2f}<br>' +
                                                 '<b>Incidentit/työntekijä:</b> %{customdata[1]:.2f}<br>' +
                                                 '<extra></extra>',
                                    customdata=hourly_stats[['avg_incidents', 'incidents_per_worker']].values
                                )
                            
                            # Yhteinen hover-tyyli kaikille kaavioille
                            fig.update_layout(
                                height=500,
                                hovermode='x unified',
                                hoverlabel=dict(
                                    bgcolor="white",
                                    font_size=14,
                                    font_family="Arial",
                                    bordercolor="gray"
                                )
                            )
                            st.plotly_chart(fig, use_container_width=True)
                            
                        except Exception as e:
                            st.error(f"Virhe kaavion luonnissa: {str(e)}")
                            st.info("Näytetään data taulukkona:")
                            st.dataframe(hourly_stats)
                    else:
                        st.warning("Ei dataa kaavion piirtämiseen.")
                
                with tab3:
                    st.subheader("📅 Kuukausinäkymä")
                    
                    if len(daily_stats) >= 1:
                        # Luo kalenterinäkymä
                        try:
                            calendar_html = create_calendar_view(daily_stats)
                            if calendar_html:
                                # Käytä korkeampaa height-arvoa jotta koko kalenteri mahtuu
                                import streamlit.components.v1 as components
                                components.html(calendar_html, height=900, scrolling=True)
                            else:
                                st.warning("Kalenterin luonti epäonnistui.")
                        except Exception as e:
                            st.error(f"Virhe kalenterin luonnissa: {str(e)}")
                            st.info("Näytetään data taulukkona:")
                            st.dataframe(daily_stats)
                        
                        # Kuukausistatistiikat
                        st.subheader("📊 Kuukauden yhteenveto")
                        col1, col2, col3, col4 = st.columns(4)
                        
                        day_target_met = len(daily_stats[daily_stats['day_target_met']]) 
                        night_target_met = len(daily_stats[daily_stats['night_target_met']])
                        total_days = len(daily_stats)
                        
                        with col1:
                            st.metric("Päivätyöntekijät", f"{day_target_met}/{total_days}", f"{day_target_met/total_days*100:.1f}%")
                        with col2:
                            st.metric("Yötyöntekijät", f"{night_target_met}/{total_days}", f"{night_target_met/total_days*100:.1f}%")
                        with col3:
                            max_day = daily_stats.loc[daily_stats['total_incidents'].idxmax()]
                            st.metric("Kiireisin päivä", f"{max_day['day']:.0f}. ({max_day['day_name']})", f"{max_day['total_incidents']:.0f} inc")
                        with col4:
                            min_day = daily_stats.loc[daily_stats['total_incidents'].idxmin()]
                            st.metric("Rauhallisin päivä", f"{min_day['day']:.0f}. ({min_day['day_name']})", f"{min_day['total_incidents']:.0f} inc")
                        
                        # Päivittäinen kehitys
                        try:
                            fig_daily = px.line(
                                daily_stats, 
                                x='date', 
                                y=['day_shift_avg', 'night_shift_avg'],
                                title='Päivittäinen kehitys',
                                labels={
                                    'value': 'Inc/työnt./h', 
                                    'variable': 'Vuoro',
                                    'date': 'Päivämäärä'
                                },
                                hover_data={
                                    'date': False,
                                    'value': ':.2f'
                                }
                            )
                            
                            # Muuta legendan nimet suomeksi ja paranna tooltip
                            fig_daily.for_each_trace(
                                lambda t: t.update(
                                    name='Päivätyöntekijät' if 'day_shift_avg' in t.name else 'Yötyöntekijät',
                                    hovertemplate='<b>Päivämäärä:</b> %{x}<br>' +
                                                 '<b>' + ('Päivätyöntekijät' if 'day_shift_avg' in t.name else 'Yötyöntekijät') + ':</b> %{y:.2f}<br>' +
                                                 '<extra></extra>'
                                )
                            )
                            
                            fig_daily.add_hline(y=5.1, line_dash="dash", line_color="red", 
                                              annotation_text="Päivätyöntekijöiden tavoite (5.1)")
                            fig_daily.add_hline(y=4.6, line_dash="dash", line_color="blue", 
                                              annotation_text="Yötyöntekijöiden tavoite (4.6)")
                            
                            fig_daily.update_layout(
                                hovermode='x unified',
                                hoverlabel=dict(
                                    bgcolor="white",
                                    font_size=14,
                                    font_family="Arial",
                                    bordercolor="gray"
                                )
                            )
                            st.plotly_chart(fig_daily, use_container_width=True)
                        except Exception as e:
                            st.error(f"Virhe päivittäisen kehityksen kaavion luonnissa: {str(e)}")
                            st.info("Näytetään data taulukkona:")
                            st.dataframe(daily_stats[['date', 'day_shift_avg', 'night_shift_avg']])
                        
                        # Päivittäinen taulukko
                        st.subheader("📋 Päivittäiset tulokset")
                        daily_display = daily_stats.copy()
                        daily_display['Päivätyöntekijät'] = daily_display.apply(
                            lambda x: f"{x['day_shift_avg']:.2f} {'✅' if x['day_target_met'] else '❌'}", axis=1
                        )
                        daily_display['Yötyöntekijät'] = daily_display.apply(
                            lambda x: f"{x['night_shift_avg']:.2f} {'✅' if x['night_target_met'] else '❌'}", axis=1
                        )
                        
                        st.dataframe(
                            daily_display[['date', 'day_name', 'total_incidents', 'Päivätyöntekijät', 'Yötyöntekijät']],
                            column_config={
                                'date': 'Päivämäärä',
                                'day_name': 'Viikonpäivä',
                                'total_incidents': 'Yhteensä inc.',
                                'Päivätyöntekijät': 'Päivätyöntekijät',
                                'Yötyöntekijät': 'Yötyöntekijät'
                            },
                            use_container_width=True
                        )
                    else:
                        st.info("Kuukausinäkymä vaatii vähintään yhden päivän dataa.")
                
                with tab4:
                    st.subheader("Tuntikohtaiset tilastot")
                    if len(hourly_stats) > 0:
                        st.dataframe(
                            hourly_stats,
                            column_config={
                                'hour_str': 'Kelloaika',
                                'avg_incidents': 'Keskim. incidentit',
                                'worker_count': 'Työntekijämäärä',
                                'incidents_per_worker': 'Inc/työnt./h',
                                'days_count': 'Päivien lukumäärä'
                            },
                            use_container_width=True
                        )
                    else:
                        st.warning("Ei tilastoja näytettäväksi.")
                
                with tab5:
                    st.subheader("💡 Optimointisuositukset")
                    
                    if len(hourly_stats) > 0:
                        # Ongelmatunnit päivätyöntekijöille
                        day_problems = hourly_stats[
                            (hourly_stats['hour'] >= 7) & 
                            (hourly_stats['hour'] < 23) & 
                            (hourly_stats['incidents_per_worker'] < 5.1)
                        ]
                        
                        # Ongelmatunnit yötyöntekijöille  
                        night_problems = hourly_stats[
                            ((hourly_stats['hour'] >= 23) | (hourly_stats['hour'] < 7)) & 
                            (hourly_stats['incidents_per_worker'] < 4.6)
                        ]
                        
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            st.markdown("### 🌅 Päivätyöntekijät")
                            if len(day_problems) > 0:
                                st.error(f"Ongelmia {len(day_problems)} tunnissa:")
                                for _, row in day_problems.iterrows():
                                    st.write(f"- {row['hour_str']}: {row['incidents_per_worker']} inc/työnt./h")
                                st.markdown("**Suositus:** Vähennä henkilöstöä ali-tuottavina aikoina tai siirrä tehtäviä.")
                            else:
                                st.success("✅ Kaikki tunnit täyttävät tavoitteen!")
                        
                        with col2:
                            st.markdown("### 🌙 Yötyöntekijät")
                            if len(night_problems) > 0:
                                st.error(f"Ongelmia {len(night_problems)} tunnissa:")
                                for _, row in night_problems.iterrows():
                                    st.write(f"- {row['hour_str']}: {row['incidents_per_worker']} inc/työnt./h")
                                st.markdown("**Suositus:** Lisää henkilöstöä ongelmallisina aikoina.")
                            else:
                                st.success("✅ Kaikki tunnit täyttävät tavoitteen!")
                        
                        # Kokonaiskuva
                        st.markdown("### 📊 Kokonaisarvio")
                        if day_avg >= 5.1 and night_avg >= 4.6:
                            st.success("🎉 Molemmat tuottavuustavoitteet saavutettu! Jatka samalla strategialla.")
                        elif day_avg >= 5.1:
                            st.warning("⚠️ Päivätyöntekijöiden tavoite saavutettu, mutta yötyöntekijät tarvitsevat parannusta.")
                        elif night_avg >= 4.6:
                            st.warning("⚠️ Yötyöntekijöiden tavoite saavutettu, mutta päivätyöntekijät tarvitsevat parannusta.")
                        else:
                            st.error("❌ Kumpikaan tuottavuustavoite ei täyty. Tarvitaan merkittäviä toimenpiteitä.")
                    else:
                        st.warning("Ei dataa suositusten tekemiseen.")
        
        except Exception as e:
            st.error(f"Virhe tiedoston käsittelyssä: {str(e)}")
            st.info("Tarkista että Excel-tiedosto sisältää sarakkeet 'Hour' ja 'Incidents handled by agent' ja että ne sisältävät numeroita.")
    
    else:
        # Ohjeet kun ei tiedostoa ladattu
        st.info("👆 Lataa Excel-tiedosto sivupalkista aloittaaksesi analyysin.")
        
        # Näytä silti PowerPoint-valinnat
        st.markdown("---")
        st.subheader("📑 PowerPoint-esityksen esikatseluvaihtoehdot")
        st.info("💡 PowerPoint-diat aktivoituvat kun olet ladannut ja käsitellyt Excel-tiedoston")
        
        # Näytä diojen selitykset
        slide_descriptions = {
            "📊 Yhteenveto": "Kokonaisanalyysi ajanjakson tuloksista, incidenttimääristä ja tavoitteiden täyttymisestä",
            "🎯 Tuottavuustavoitteet": "Yksityiskohtainen analyysi päivä- ja yötyöntekijöiden tuottavuustavoitteista värillisine tuloksineen",
            "📈 Yhdistetty kaavio": "Visuaalinen esitys tuntikohtaisesta kuormituksesta ja henkilöstömääristä",
            "🕐 Tuntikohtainen analyysi": "Syvällinen tarkastelu huipputunneista, vuorojärjestelystä ja tehokkuudesta",
            "📅 Kuukausinäkymä": "Päivittäiset trendit, kuukauden yhteenveto ja kehityksen seuranta",
            "💡 Suositukset": "Konkreettiset toimenpide-ehdotukset optimointiin ja ongelmien ratkaisuun"
        }
        
        for title, description in slide_descriptions.items():
            with st.expander(title):
                st.write(description)
        
        st.markdown("---")
        st.subheader("📋 Käyttöohjeet")
        st.markdown("""
        1. **Lataa Excel-tiedosto** sivupalkista
        2. Tiedoston tulee sisältää vähintään sarakkeet:
           - `Hour` (0-23, numeroina)
           - `Incidents handled by agent` (määrä, numeroina)
           - `Date` (valinnainen, päivämäärille)
        3. **Valitse PowerPoint-diat** sivupalkista
        4. **Tarkastele tuloksia** eri välilehdiltä:
           - 📊 Yhdistetty näkymä
           - 📈 Tuntikohtainen analyysi  
           - 📅 Kuukausinäkymä
           - 📋 Tilastot
           - 💡 Suositukset
        5. **Lataa PowerPoint-esitys** valituista dioista
        """)
        
        st.markdown("### 🎯 Mitä työkalu analysoi:")
        st.markdown("""
        - **Tuottavuustavoitteiden täyttyminen** vuoroittain
        - **Tuntikohtaiset kuormitukset** ja henkilöstötarpeet
        - **Päivittäiset suorituskykytrendit** 
        - **Optimointisuositukset** resurssien allokointiin
        - **Interaktiiviset visualisoinnit** helposti ymmärrettävässä muodossa
        - **PowerPoint-esitykset** räätälöityjen diojen kanssa
        """)

        # Näytä esimerkki oikeasta datamuodosta
        st.markdown("### 📝 Esimerkki oikeasta datamuodosta:")
        example_data = pd.DataFrame({
            'Date': ['2025-02-01', '2025-02-01', '2025-02-01'],
            'Hour': [0, 1, 2],
            'Incidents handled by agent': [9, 14, 16]
        })
        st.dataframe(example_data, use_container_width=True)
        
        # PowerPoint-ominaisuuksien esittely
        st.markdown("### 📑 PowerPoint-ominaisuudet:")
        st.markdown("""
        **Saatavilla olevat diat:**
        - 📊 **Yhteenveto**: Kokonaisanalyysi ja päätulokset
        - 🎯 **Tuottavuustavoitteet**: Yksityiskohtainen tavoiteanalyysi
        - 📊 **Yhdistetty kaavio**: Visuaalinen tuntikohtainen analyysi
        - 📈 **Tuntikohtainen analyysi**: Syvällinen tuntitason tarkastelu
        - 📅 **Kuukausinäkymä**: Päivittäiset trendit ja yhteenvedot
        - 💡 **Suositukset**: Toimenpide-ehdotukset ja optimointi
        
        **Kuinka käyttää:**
        1. Valitse haluamasi diat sivupalkista
        2. Klikkaa "Luo PowerPoint" -painiketta
        3. Lataa valmis esitys suoraan selaimesta
        """)

if __name__ == "__main__":
    main()
